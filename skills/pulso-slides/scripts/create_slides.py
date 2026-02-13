#!/usr/bin/env python3
"""
Create minimal PowerPoint presentations from Markdown input.

Usage:
    python create_slides.py <input.md> [output.pptx]
    echo "# Title" | python create_slides.py -

Markdown format:
    # Title          -> slide title
    ## Subtitle      -> subtitle (title slide only)
    - Bullet item    -> bulleted line
    Plain text       -> paragraph
    ---              -> slide separator
"""

import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn, nsmap
from lxml import etree


# -- Layout constants (inches) ------------------------------------------------

SLIDE_W, SLIDE_H = 13.333, 7.5

TITLE_SLIDE = {
    "title":    {"left": 1.0, "top": 2.5, "width": 11.333, "height": 1.5},
    "subtitle": {"left": 2.0, "top": 4.2, "width": 9.333,  "height": 1.0},
}

CONTENT_SLIDE = {
    "title": {"left": 0.75, "top": 0.5,  "width": 11.833, "height": 0.8},
    "body":  {"left": 1.0,  "top": 1.8,  "width": 11.333, "height": 5.2},
}

FONT_FAMILY = "Open Sans"


# -- Markdown parser -----------------------------------------------------------

def parse_markdown(text: str) -> list[dict]:
    """Parse markdown text into a list of slide dicts.

    Each slide dict has:
        kind: "title" | "content"
        title: str
        subtitle: str | None       (title slides only)
        body: list of (type, text)  where type is "bullet" or "paragraph"
    """
    raw_slides = text.strip().split("\n---\n")
    slides = []

    for raw in raw_slides:
        lines = raw.strip().splitlines()
        if not lines:
            continue

        title = ""
        subtitle = None
        body: list[tuple[str, str]] = []

        for line in lines:
            stripped = line.strip()
            if stripped.startswith("# "):
                title = stripped[2:].strip()
            elif stripped.startswith("## "):
                subtitle = stripped[3:].strip()
            elif stripped.startswith("- "):
                body.append(("bullet", stripped[2:].strip()))
            elif stripped:
                body.append(("paragraph", stripped))

        # Determine slide kind: title slide if it has a subtitle, or if it
        # has only a title with no body.
        is_title = subtitle is not None or (title and not body)
        slides.append({
            "kind": "title" if is_title else "content",
            "title": title,
            "subtitle": subtitle,
            "body": body,
        })

    return slides


# -- PPTX builder -------------------------------------------------------------

def _add_textbox(slide, left, top, width, height):
    """Add a textbox and return its text_frame."""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    return tf


def _set_font(run, size_pt, bold=False, color_hex="000000"):
    """Apply font settings to a run."""
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.name = FONT_FAMILY
    run.font.color.rgb = _hex_to_rgb(color_hex)


def _hex_to_rgb(hex_str):
    from pptx.dml.color import RGBColor
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _set_bullet(paragraph):
    """Set a middle-dot bullet on a paragraph via OOXML XML manipulation.

    Produces proper <a:buChar char="·"/> with hanging indent so PowerPoint
    treats it as a real editable bullet.
    """
    pPr = paragraph._p.get_or_add_pPr()

    # Hanging indent: marL = left margin, indent = (negative) first-line indent
    pPr.set("marL", str(Emu(Inches(0.35))))
    pPr.set("indent", str(Emu(-Inches(0.25))))

    # Remove any existing bullet definitions
    for tag in ("a:buNone", "a:buChar", "a:buAutoNum"):
        for old in pPr.findall(qn(tag)):
            pPr.remove(old)

    # Add buChar element with middle dot
    buChar = etree.SubElement(pPr, qn("a:buChar"))
    buChar.set("char", "\u00B7")

    # Set bullet font to match body font
    buFont = etree.SubElement(pPr, qn("a:buFont"))
    buFont.set("typeface", FONT_FAMILY)


def _set_spacing(paragraph, before_pt=0, after_pt=0):
    """Set paragraph spacing via XML."""
    pPr = paragraph._p.get_or_add_pPr()
    spcBef = etree.SubElement(pPr, qn("a:spcBef"))
    spcPts = etree.SubElement(spcBef, qn("a:spcPts"))
    spcPts.set("val", str(int(before_pt * 100)))

    spcAft = etree.SubElement(pPr, qn("a:spcAft"))
    spcPts2 = etree.SubElement(spcAft, qn("a:spcPts"))
    spcPts2.set("val", str(int(after_pt * 100)))


def _build_title_slide(prs, slide_data):
    """Build a title slide (centered title + optional subtitle)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Title
    cfg = TITLE_SLIDE["title"]
    tf = _add_textbox(slide, **cfg)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = slide_data["title"]
    _set_font(run, 36, bold=True, color_hex="000000")

    # Vertically center title in its box
    bodyPr = tf._txBody.find(qn("a:bodyPr"))
    bodyPr.set("anchor", "ctr")

    # Subtitle
    if slide_data["subtitle"]:
        cfg = TITLE_SLIDE["subtitle"]
        tf2 = _add_textbox(slide, **cfg)
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = slide_data["subtitle"]
        _set_font(run2, 20, bold=False, color_hex="444444")


def _build_content_slide(prs, slide_data):
    """Build a content slide (top-left title + body with bullets/paragraphs)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Title
    cfg = CONTENT_SLIDE["title"]
    tf = _add_textbox(slide, **cfg)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = slide_data["title"]
    _set_font(run, 28, bold=True, color_hex="000000")

    # Vertically center the title in its box
    bodyPr = tf._txBody.find(qn("a:bodyPr"))
    bodyPr.set("anchor", "ctr")

    # Body
    if not slide_data["body"]:
        return

    cfg = CONTENT_SLIDE["body"]
    tf_body = _add_textbox(slide, **cfg)

    first = True
    for kind, text in slide_data["body"]:
        if first:
            p = tf_body.paragraphs[0]
            first = False
        else:
            p = tf_body.add_paragraph()

        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = text
        _set_font(run, 18, bold=False, color_hex="1A1A1A")

        if kind == "bullet":
            _set_bullet(p)
            _set_spacing(p, before_pt=4, after_pt=4)
        else:
            _set_spacing(p, before_pt=6, after_pt=6)


def build_presentation(slides_data: list[dict], output_path: str):
    """Create the PPTX file from parsed slide data."""
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    for slide_data in slides_data:
        if slide_data["kind"] == "title":
            _build_title_slide(prs, slide_data)
        else:
            _build_content_slide(prs, slide_data)

    prs.save(output_path)
    print(f"Saved: {output_path} ({len(slides_data)} slides)")


# -- CLI -----------------------------------------------------------------------

def main():
    if len(sys.argv) < 2:
        print(__doc__.strip())
        sys.exit(1)

    input_arg = sys.argv[1]

    # Read markdown from stdin or file
    if input_arg == "-":
        md_text = sys.stdin.read()
        output_path = sys.argv[2] if len(sys.argv) > 2 else "output.pptx"
    else:
        input_path = Path(input_arg)
        if not input_path.exists():
            print(f"Error: file not found: {input_path}")
            sys.exit(1)
        md_text = input_path.read_text(encoding="utf-8")
        if len(sys.argv) > 2:
            output_path = sys.argv[2]
        else:
            output_path = str(input_path.with_suffix(".pptx"))

    slides_data = parse_markdown(md_text)
    if not slides_data:
        print("Error: no slides found in input")
        sys.exit(1)

    build_presentation(slides_data, output_path)


if __name__ == "__main__":
    main()
